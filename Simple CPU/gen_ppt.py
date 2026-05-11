import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Precise HTML Colors
    c_bg_page = RGBColor(248, 250, 252)         # #f8fafc
    c_card = RGBColor(255, 255, 255)            # #ffffff
    c_text_main = RGBColor(15, 23, 42)          # #0f172a
    c_text_muted = RGBColor(51, 65, 85)         # #334155
    c_primary = RGBColor(37, 99, 235)           # #2563eb
    c_secondary = RGBColor(5, 150, 105)         # #059669
    c_border = RGBColor(203, 213, 225)          # #cbd5e1
    c_code_bg = RGBColor(30, 41, 59)            # #1e293b
    c_code_text = RGBColor(248, 250, 252)       # #f8fafc
    c_code_border = RGBColor(71, 85, 105)       # #475569
    c_hl_bg = RGBColor(239, 246, 255)           # #eff6ff

    font_main = "Segoe UI"
    font_code = "Consolas"

    def add_base(slide, current, total):
        # 1. Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = c_bg_page
        
        # 2. Main White Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.5), Inches(0.375), 
            Inches(12.333), Inches(6.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = c_card
        card.line.color.rgb = c_border
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.02 # smaller border radius
        
        # Add shadow to card for that "web" feel
        spPr = card.element.spPr
        effectLst = SubElement(spPr, 'a:effectLst')
        outerShdw = SubElement(effectLst, 'a:outerShdw', blurRad="150000", dist="50000", dir="5400000", algn="b")
        srgbClr = SubElement(outerShdw, 'a:srgbClr', val="000000")
        SubElement(srgbClr, 'a:alpha', val="10000") # 10% opacity
        
        card.text_frame.clear()
        
        # 3. Slide counter
        counter = slide.shapes.add_textbox(Inches(11.5), Inches(6.5), Inches(1.0), Inches(0.5))
        p = counter.text_frame.paragraphs[0]
        p.text = f"{current} / {total}"
        p.font.size = Pt(14)
        p.font.name = font_main
        p.font.color.rgb = c_text_muted
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT

    def add_header(slide, title, badge_text):
        # Title
        txBox = slide.shapes.add_textbox(Inches(1.25), Inches(0.8), Inches(9.0), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.name = font_main
        p.font.color.rgb = c_text_main
        
        # Badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(10.5), Inches(0.9), Inches(1.5), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = c_primary
        badge.line.fill.background()
        badge.adjustments[0] = 0.2
        p_badge = badge.text_frame.paragraphs[0]
        p_badge.text = badge_text
        p_badge.font.size = Pt(12)
        p_badge.font.name = font_code
        p_badge.font.bold = True
        p_badge.alignment = PP_ALIGN.CENTER
        
        # Blue Line (perfectly spans content area 1.25 to 12.083)
        line = slide.shapes.add_connector(1, Inches(1.25), Inches(1.5), Inches(12.083), Inches(1.5))
        line.line.color.rgb = c_primary
        line.line.width = Pt(3)

    def add_text(slide, text, left, top, width, height=0.5, font_size=16, color=c_text_muted, bold=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_main
        p.font.color.rgb = color
        p.font.bold = bold
        return txBox

    def add_bullets(slide, items, left, top, width, height=3.0):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "■  " + item
            p.font.size = Pt(16)
            p.font.name = font_main
            p.font.color.rgb = c_text_muted
            # Small paragraph spacing
            p.space_after = Pt(10)

    def add_highlight_box(slide, title, content, left, top, width, height):
        # Light blue box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_hl_bg
        shape.line.fill.background()
        shape.adjustments[0] = 0.05
        
        # Primary left border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(left), Inches(top), Inches(0.08), Inches(height)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = c_primary
        border.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.name = font_main
        p.font.color.rgb = c_primary
        p.space_after = Pt(5)
        
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.name = font_main
        p2.font.color.rgb = c_text_main

    def add_code(slide, code_str, left, top, width, height):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_code_bg
        shape.line.color.rgb = c_code_border
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.05
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = code_str
        p.font.size = Pt(12)
        p.font.name = font_code
        p.font.color.rgb = c_code_text

    total = 11

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 1, total)
    add_text(slide, "Simple CPU Emulator", 1.25, 2.0, 10.833, font_size=50, color=c_primary, bold=True).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, "Deep Dive into C Implementation", 1.25, 3.0, 10.833, font_size=24, color=c_text_main, bold=True).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, "A highly detailed breakdown of the exact logic running inside our core source files (`.c`), detailing hardware simulation mechanisms and memory manipulation.", 2.5, 3.8, 8.333, font_size=16, color=c_text_muted).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2. Project Structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 2, total)
    add_header(slide, "Project Implementation Structure", "Core Files")
    add_text(slide, "The executable logic is divided into four distinct C source files. We will strictly focus on how these files operate natively.", 1.25, 1.8, 10.833)
    add_highlight_box(slide, "1. registers.c", "Allocates the simulated memory array, creates global variables representing hardware registers, and handles masked write limits.", 1.25, 2.8, 5.1, 1.4)
    add_highlight_box(slide, "2. fetchDecodeExecute.c", "Implements the core machine loop functions. Manages program counter incrementing and opcode isolation via bitwise operations.", 1.25, 4.4, 5.1, 1.4)
    add_highlight_box(slide, "3. instructions.c", "Contains the specific algorithmic definitions for the 15+ supported instructions, memory pointer handling, and ASCII visualization.", 6.983, 2.8, 5.1, 1.4)
    add_highlight_box(slide, "4. assembler.c", "Provides the sscanf based string-parsing engine that builds 16-bit binary integers out of human-written text.", 6.983, 4.4, 5.1, 1.4)

    # 3. Hardware State
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 3, total)
    add_header(slide, "Hardware State Simulation", "registers.c")
    add_text(slide, "This file allocates the state variables used by the entire processor. Standard integer types simulate size-specific hardware constraints.", 1.25, 1.8, 10.833)
    add_bullets(slide, [
        "Memory: Simulated as a massive array: int memory[4096]",
        "16-bit Registers: AC (Accumulator), IR (Instruction Register), MBR",
        "12-bit Registers: PC (Program Counter) and MAR",
        "8-bit Registers: InREG and OutREG for standard I/O processing"
    ], 1.25, 2.8, 5.1)
    add_code(slide, "// registers.c globals\nint cpuClock = 0;\n\nint AC = 0, IR = 0, MBR = 0;\nint PC = 0, MAR = 0;\nint InREG = 0, OutREG = 0;\n\nint data_bus = 0;\nint memory[4096];", 6.983, 2.8, 5.1, 3.5)

    # 4. Enforcing Limits
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 4, total)
    add_header(slide, "Enforcing Hardware Limitations", "registers.c")
    add_text(slide, "C integer variables are usually 32-bits, but our CPU must behave strictly according to 8/12/16-bit boundaries. This file uses MASK macros (& 0xFFFF) to forcibly truncate.", 1.25, 1.8, 10.833)
    add_highlight_box(slide, "The CPU Clock", "Every single setter function (like write_AC) and memory interaction (like mem_read) inherently increments the global cpuClock variable, acting as an accurate performance profiler.", 1.25, 3.0, 5.1, 2.0)
    add_code(slide, "// C macros force bit limits\n#define MASK16(val) ((val)&0xFFFF)\n#define MASK12(val) ((val)&0xFFF)\n\nvoid write_AC(int val) {\n    AC = MASK16(val);\n    cpuClock++;\n}\n\nvoid write_MAR(int val) {\n    MAR = MASK12(val);\n    cpuClock++;\n}", 6.983, 2.8, 5.1, 3.8)

    # 5. Fetch Decode
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 5, total)
    add_header(slide, "The FDE Cycle: Fetch & Decode", "fetchDecodeExecute.c")
    add_text(slide, "The cycle engine lives here. The fetch() function dictates how exactly memory turns into an executable command.", 1.25, 1.8, 10.833)
    add_bullets(slide, ["Push current PC to MAR", "Read memory via data_bus", "Store result in IR", "Increment PC"], 1.25, 2.6, 5.1, 1.8)
    add_text(slide, "Next, decode() isolates the 4-bit opcode and 12-bit address variables from the 16-bit IR.", 1.25, 4.8, 5.1, 1.0)
    add_code(slide, "void fetch(void) {\n    write_MAR(PC);\n    data_bus = mem_read(MAR);\n    write_IR(data_bus);\n    write_PC(PC + 1);\n}\n\nvoid decode(void) {\n    // Shift away lower 12 bits\n    opcode = (IR >> 12) & 0xF;\n    // Mask away top 4 bits\n    write_MAR(IR & 0x0FFF);\n}", 6.983, 2.6, 5.1, 3.8)

    # 6. Execute
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 6, total)
    add_header(slide, "The FDE Cycle: Execution Router", "fetchDecodeExecute.c")
    add_text(slide, "Once the opcode integer is parsed (0x1 through 0xD), the execute() function acts as a gigantic hardware multiplexer, triggering functions exported from instructions.c.", 1.25, 1.8, 10.833)
    add_code(slide, "void execute(void) {\n    switch (opcode) {\n        case 0x1: LOAD();     break;\n        case 0x2: STORE();    break;\n        case 0x3: ADD();      break;\n        case 0x4: SUBT();     break;\n        case 0x5: IN();       break;\n        case 0x6: OUT();      break;\n        case 0x8: SKIPCOND(); break;\n        case 0xA: CLEAR();    break;\n        // ... Covers all 13 hex mappings\n    }\n}", 4.166, 2.8, 5.0, 3.8)

    # 7. Direct Logic
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 7, total)
    add_header(slide, "Direct Logic Operations", "instructions.c")
    add_text(slide, "This implementation file defines the actual manipulation of registers. No logic is abstracted; the internal data_bus is explicitly simulated.", 1.25, 1.8, 10.833)
    add_highlight_box(slide, "SKIPCOND Logic", "Instead of dedicated conditional jumps, SKIPCOND peeks at the top bits of the address field (stored in MAR) to compare the AC against zero (0x000 for <0, 0x400 for ==0), skipping the PC ahead if true.", 1.25, 2.8, 5.1, 2.0)
    add_code(slide, "void ADD(void) {\n    data_bus = mem_read(MAR);\n    write_MBR(data_bus);\n    data_bus = MBR;\n    write_AC(AC + data_bus);\n}\n\nvoid SKIPCOND(void) {\n    if      (MAR == 0x000 && AC < 0)  \n        write_PC(PC+1);\n    else if (MAR == 0x400 && AC == 0) \n        write_PC(PC+1);\n    else if (MAR == 0x800 && AC > 0)  \n        write_PC(PC+1);\n}", 6.983, 2.8, 5.1, 3.8)

    # 8. Indirect
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 8, total)
    add_header(slide, "Indirect Addressing", "instructions.c")
    add_text(slide, "In addition to standard instructions, this file manages 'Indirect' pointer-following instructions (like ADDI or LOADI) and console visualization.", 1.25, 1.8, 10.833)
    add_highlight_box(slide, "Following Pointers", "In ADDI, the MAR initially holds an address that points to another address. The code intentionally calls mem_read(MAR) twice, updating the MAR with the pointer value in between reads.", 1.25, 2.8, 5.1, 2.0)
    add_code(slide, "void ADDI(void) {\n    // Fetch pointer address\n    data_bus = mem_read(MAR);\n    \n    // Follow the pointer\n    write_MAR(data_bus);          \n    \n    // Read real value & add\n    data_bus = mem_read(MAR);\n    write_MBR(data_bus);\n    write_AC(MASK16(AC + MBR));\n}", 6.983, 2.8, 5.1, 3.8)

    # 9. Assembler
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 9, total)
    add_header(slide, "Text to Machine Code Translation", "assembler.c")
    add_text(slide, "To create runnable programs, assembler.c reads textual instructions (like \"ADD 20\") and constructs the explicit 16-bit binary equivalents.", 1.25, 1.8, 10.833)
    add_bullets(slide, [
        "sscanf splits the string into mnemonic string and operand integer.",
        "strcmp compares the string to known opcodes.",
        "It left-shifts the matched opcode by 12 bits (e.g., 0x3 << 12).",
        "It bitwise-ORs the operand onto the bottom 12 bits."
    ], 1.25, 2.8, 5.1)
    add_code(slide, "int assemble(char *line) {\n    char mnemonic[10];\n    int operand = 0;\n    \n    sscanf(line, \"%s %d\", mnemonic, &operand);\n\n    if (strcmp(mnemonic, \"LOAD\") == 0) \n        return MASK16((0x1 << 12) | MASK12(operand)); \n        \n    else if (strcmp(mnemonic, \"ADD\") == 0) \n        return MASK16((0x3 << 12) | MASK12(operand));\n    //... Handles all formats\n}", 6.983, 2.8, 5.1, 3.8)

    # 10. Test Case
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 10, total)
    add_header(slide, "Test Case Walkthrough", "program.txt")
    add_text(slide, "Let's trace how the emulator executes the simple addition program loaded from program.txt.", 1.25, 1.8, 10.833)
    add_code(slide, "; Addr   Instruction\n0x000    INPUT       ; User -> AC\n0x001    STORE 20    ; AC -> Mem[20]\n0x002    INPUT       ; User -> AC\n0x003    ADD   20    ; AC = AC + Mem[20]\n0x004    OUTPUT      ; Print AC\n0x005    HALT        ; Stop", 1.25, 2.8, 5.1, 3.8)
    add_highlight_box(slide, "Execution Trace", 
        "1. 0x000: IN() pauses. User types 15. AC = 15.\n"
        "2. 0x001: STORE() routes 15 from AC via MBR into memory[20].\n"
        "3. 0x002: IN() pauses. User types 25. AC = 25.\n"
        "4. 0x003: ADD() reads 15 from memory[20], updates AC = 25 + 15 = 40.\n"
        "5. 0x004: OUT() routes 40 to OutREG and prints to terminal.", 6.983, 2.8, 5.1, 3.8)

    # 11. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, 11, total)
    add_text(slide, "Summary", 1.25, 2.5, 10.833, font_size=50, color=c_text_main, bold=True).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(slide, "The power of this CPU emulator relies entirely on its distinct C implementation logic: registers.c enforces hardware size limits, fetchDecodeExecute.c drives the system clock and opcode isolation, instructions.c manipulates exact memory paths, and assembler.c acts as the translation layer.", 2.5, 3.8, 8.333, font_size=18, color=c_text_muted).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bottom box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.166), Inches(5.5), Inches(5.0), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c_bg_page
    shape.line.color.rgb = c_primary
    shape.line.width = Pt(3)
    p = shape.text_frame.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(24)
    p.font.name = font_main
    p.font.bold = True
    p.font.color.rgb = c_primary

    prs.save('SimpleCPU_Presentation_v2.pptx')

if __name__ == '__main__':
    create_presentation()
